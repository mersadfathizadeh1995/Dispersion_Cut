"""Main PublicationFigureDialog class."""

from __future__ import annotations

from typing import Optional, Dict, Any, List, Tuple
from pathlib import Path
from matplotlib.backends import qt_compat

QtWidgets = qt_compat.QtWidgets
QtGui = qt_compat.QtGui
QtCore = qt_compat.QtCore

from ..pub_figures import PublicationFigureGenerator, PlotConfig
from .qt_compat import (
    _get_qt_orientation_horizontal,
    _get_qt_align_top,
    _get_qt_user_role,
    _get_qt_item_is_selectable,
    _get_qt_item_is_enabled,
    _get_qt_item_is_user_checkable,
    _get_qt_checked,
    _get_qt_unchecked,
    _get_qt_extended_selection,
    _get_qt_msgbox_yes,
    _get_qt_msgbox_no,
)
from .constants import FIGURE_TYPES


class PublicationFigureDialog(QtWidgets.QDialog):
    """Dialog for configuring and generating publication-quality figures."""

    def __init__(self, controller, parent: Optional[QtWidgets.QWidget] = None):
        super().__init__(parent)
        self.controller = controller
        self.setWindowTitle("Export Publication Figure")
        self.resize(600, 700)

        # Main layout
        layout = QtWidgets.QVBoxLayout(self)

        # Tab widget for organized sections
        self.tabs = QtWidgets.QTabWidget(self)
        layout.addWidget(self.tabs)

        # Create tabs
        self._build_plot_type_tab()
        self._build_styling_tab()
        self._build_axes_tab()
        self._build_output_tab()

        # Buttons
        button_box = QtWidgets.QDialogButtonBox(self)
        try:
            generate_btn = QtWidgets.QDialogButtonBox.Ok
            cancel_btn = QtWidgets.QDialogButtonBox.Cancel
        except AttributeError:
            generate_btn = QtWidgets.QDialogButtonBox.StandardButton.Ok
            cancel_btn = QtWidgets.QDialogButtonBox.StandardButton.Cancel

        button_box.setStandardButtons(generate_btn | cancel_btn)
        # Rename OK button to "Generate"
        try:
            ok_button = button_box.button(generate_btn)
            if ok_button:
                ok_button.setText("Generate")
        except Exception:
            pass

        button_box.accepted.connect(self._on_generate)
        button_box.rejected.connect(self.reject)

        layout.addWidget(button_box)

    def _build_plot_type_tab(self):
        """Build the Plot Type selection tab with tree view."""
        tab = QtWidgets.QWidget()
        layout = QtWidgets.QVBoxLayout(tab)

        # Search box at the top
        search_layout = QtWidgets.QHBoxLayout()
        search_label = QtWidgets.QLabel("Search:")
        self.plot_search_edit = QtWidgets.QLineEdit()
        self.plot_search_edit.setPlaceholderText("Type to filter figure types...")
        self.plot_search_edit.setClearButtonEnabled(True)
        self.plot_search_edit.textChanged.connect(self._filter_plot_types)
        search_layout.addWidget(search_label)
        search_layout.addWidget(self.plot_search_edit)
        layout.addLayout(search_layout)

        # Splitter for tree view (left) and description panel (right)
        splitter = QtWidgets.QSplitter(_get_qt_orientation_horizontal())

        # Left panel: Tree view of figure types
        left_widget = QtWidgets.QWidget()
        left_layout = QtWidgets.QVBoxLayout(left_widget)
        left_layout.setContentsMargins(0, 0, 0, 0)

        tree_label = QtWidgets.QLabel("<b>Figure Types (check to select)</b>")
        left_layout.addWidget(tree_label)

        self.plot_tree = QtWidgets.QTreeWidget()
        self.plot_tree.setHeaderHidden(True)
        self.plot_tree.setRootIsDecorated(True)
        # Use single selection for clicking, but checkboxes for multi-select
        self.plot_tree.setSelectionMode(QtWidgets.QAbstractItemView.SingleSelection
                                        if hasattr(QtWidgets.QAbstractItemView, 'SingleSelection')
                                        else QtWidgets.QAbstractItemView.SelectionMode.SingleSelection)

        left_layout.addWidget(self.plot_tree)
        splitter.addWidget(left_widget)

        # Right panel: Description (create BEFORE connecting signals and populating tree)
        right_widget = QtWidgets.QWidget()
        right_layout = QtWidgets.QVBoxLayout(right_widget)
        right_layout.setContentsMargins(0, 0, 0, 0)

        desc_title_label = QtWidgets.QLabel("<b>Description</b>")
        right_layout.addWidget(desc_title_label)

        self.plot_desc_label = QtWidgets.QLabel("Select a figure type to see its description.")
        self.plot_desc_label.setWordWrap(True)
        self.plot_desc_label.setAlignment(_get_qt_align_top())
        self.plot_desc_label.setStyleSheet("padding: 10px; background-color: palette(base);")
        self.plot_desc_label.setMinimumHeight(100)
        right_layout.addWidget(self.plot_desc_label)

        # Status indicator
        self.plot_status_label = QtWidgets.QLabel("")
        self.plot_status_label.setWordWrap(True)
        right_layout.addWidget(self.plot_status_label)

        # Selection count indicator
        self.selection_count_label = QtWidgets.QLabel("")
        self.selection_count_label.setStyleSheet("font-weight: bold; color: #0066cc;")
        right_layout.addWidget(self.selection_count_label)

        right_layout.addStretch()
        splitter.addWidget(right_widget)

        # Set initial splitter sizes (60% tree, 40% description)
        splitter.setSizes([300, 200])

        layout.addWidget(splitter, 1)  # Give splitter stretch priority

        # Per-offset options (create BEFORE populating tree since selection may trigger update)
        offset_group = QtWidgets.QGroupBox("Offset Options")
        offset_layout = QtWidgets.QFormLayout(offset_group)

        # For per-offset plot types (existing functionality)
        self.max_offsets_spin = QtWidgets.QSpinBox()
        self.max_offsets_spin.setRange(1, 50)
        self.max_offsets_spin.setValue(10)
        offset_layout.addRow("Maximum offsets to plot:", self.max_offsets_spin)

        # For individual offset analysis (single offset selection)
        self.offset_selector_combo = QtWidgets.QComboBox()
        self.offset_selector_combo.setMinimumWidth(200)
        self.offset_selector_label = QtWidgets.QLabel("Select offset:")
        offset_layout.addRow(self.offset_selector_label, self.offset_selector_combo)

        # Populate offset selector from controller
        self._populate_offset_selector()

        # ------------------------------------------------------------------
        # Grid offset selection (for Comparison Grid)
        # ------------------------------------------------------------------
        self.grid_offset_selection_label = QtWidgets.QLabel("<b>Select Offsets for Grid:</b>")
        offset_layout.addRow(self.grid_offset_selection_label)
        
        # Selection buttons
        grid_buttons_layout = QtWidgets.QHBoxLayout()
        self.select_all_btn = QtWidgets.QPushButton("Select All")
        self.deselect_all_btn = QtWidgets.QPushButton("Deselect All")
        self.select_all_btn.clicked.connect(self._select_all_grid_offsets)
        self.deselect_all_btn.clicked.connect(self._deselect_all_grid_offsets)
        grid_buttons_layout.addWidget(self.select_all_btn)
        grid_buttons_layout.addWidget(self.deselect_all_btn)
        grid_buttons_layout.addStretch()
        offset_layout.addRow("", grid_buttons_layout)
        
        # List widget with checkboxes for offset selection
        self.grid_offset_list = QtWidgets.QListWidget()
        self.grid_offset_list.setMinimumHeight(120)
        self.grid_offset_list.setMaximumHeight(180)
        self.grid_offset_list.setSelectionMode(QtWidgets.QAbstractItemView.SelectionMode.NoSelection)
        offset_layout.addRow("", self.grid_offset_list)
        
        # Populate the grid offset list
        self._populate_grid_offset_list()

        # ------------------------------------------------------------------
        # Grid display mode options (for offset_grid)
        # ------------------------------------------------------------------
        grid_display_label = QtWidgets.QLabel("Grid display mode:")
        self.grid_mode_curves = QtWidgets.QRadioButton("Curves Only")
        self.grid_mode_spectrum = QtWidgets.QRadioButton("Spectrum Only")
        self.grid_mode_both = QtWidgets.QRadioButton("Curves + Spectrum")
        self.grid_mode_curves.setChecked(True)  # Default to curves only
        
        grid_mode_layout = QtWidgets.QHBoxLayout()
        grid_mode_layout.addWidget(self.grid_mode_curves)
        grid_mode_layout.addWidget(self.grid_mode_spectrum)
        grid_mode_layout.addWidget(self.grid_mode_both)
        grid_mode_layout.addStretch()
        offset_layout.addRow(grid_display_label, grid_mode_layout)

        # Grid layout options for offset_grid
        grid_options_layout = QtWidgets.QHBoxLayout()
        grid_options_layout.addWidget(QtWidgets.QLabel("Grid rows:"))
        self.grid_rows_spin = QtWidgets.QSpinBox()
        self.grid_rows_spin.setRange(1, 10)
        self.grid_rows_spin.setValue(0)  # 0 = auto
        self.grid_rows_spin.setSpecialValueText("Auto")
        grid_options_layout.addWidget(self.grid_rows_spin)
        grid_options_layout.addWidget(QtWidgets.QLabel("cols:"))
        self.grid_cols_spin = QtWidgets.QSpinBox()
        self.grid_cols_spin.setRange(1, 10)
        self.grid_cols_spin.setValue(0)  # 0 = auto
        self.grid_cols_spin.setSpecialValueText("Auto")
        grid_options_layout.addWidget(self.grid_cols_spin)
        grid_options_layout.addStretch()
        offset_layout.addRow("", grid_options_layout)

        offset_note = QtWidgets.QLabel("Note: Colorbar settings are in Styling tab > Spectrum Options")
        offset_note.setStyleSheet("color: gray; font-style: italic;")
        offset_layout.addRow("", offset_note)

        layout.addWidget(offset_group)
        self.offset_group = offset_group
        self.offset_group.setEnabled(False)

        # NOW connect signals and populate tree (after all widgets exist)
        self.plot_tree.itemSelectionChanged.connect(self._on_tree_selection_changed)
        self.plot_tree.itemChanged.connect(self._on_tree_item_changed)
        self._populate_plot_tree()

        # Bottom options - removed obsolete checkbox, selection is via checkboxes in tree
        # The tree checkboxes now handle multi-selection

        self.tabs.addTab(tab, "Plot Type")

    def _populate_plot_tree(self):
        """Populate the tree widget with figure type categories."""
        self.plot_tree.clear()

        first_implemented_item = None
        
        # Get available offsets for Source Offset Analysis sub-items
        available_offsets = self._get_available_offsets()

        for category_name, figure_types in FIGURE_TYPES.items():
            # Create category item
            category_item = QtWidgets.QTreeWidgetItem([category_name])
            category_item.setFlags(category_item.flags() & ~_get_qt_item_is_selectable())
            category_item.setExpanded(True)

            # Make category header bold
            font = category_item.font(0)
            font.setBold(True)
            category_item.setFont(0, font)

            for display_name, internal_key, description, is_implemented in figure_types:
                # Create figure type item
                if is_implemented:
                    item_text = display_name
                else:
                    item_text = f"{display_name} (Coming Soon)"

                fig_item = QtWidgets.QTreeWidgetItem([item_text])
                fig_item.setData(0, _get_qt_user_role(), {
                    'key': internal_key,
                    'name': display_name,
                    'description': description,
                    'implemented': is_implemented
                })

                # For individual offset types, add offset sub-items instead of checkboxes
                if internal_key in ['offset_curve_only', 'offset_with_spectrum', 'offset_spectrum_only']:
                    # These types have offset children - make parent non-checkable
                    fig_item.setFlags(fig_item.flags() & ~_get_qt_item_is_user_checkable())
                    
                    # Determine label suffix and filter based on type
                    if internal_key == 'offset_curve_only':
                        label_suffix = "[curve]"
                        filter_with_spectrum = False  # Show all offsets
                    elif internal_key == 'offset_with_spectrum':
                        label_suffix = "[curve + spectrum]"
                        filter_with_spectrum = True  # Only offsets with spectrum
                    else:  # offset_spectrum_only
                        label_suffix = "[spectrum]"
                        filter_with_spectrum = True  # Only offsets with spectrum
                    
                    # Add sub-items for each offset
                    offset_added = False
                    for offset_info in available_offsets:
                        # Filter based on spectrum availability
                        if filter_with_spectrum and not offset_info.get('has_spectrum', False):
                            continue
                        
                        offset_text = f"{offset_info['name']} {label_suffix}"
                        offset_item = QtWidgets.QTreeWidgetItem([offset_text])
                        offset_item.setData(0, _get_qt_user_role(), {
                            'key': internal_key,
                            'name': offset_info['name'],
                            'description': f"{description}\n\nOffset: {offset_info['name']}",
                            'implemented': True,
                            'offset_index': offset_info['index'],
                            'offset_name': offset_info['name'],
                            'has_spectrum': offset_info.get('has_spectrum', False)
                        })
                        offset_item.setFlags(offset_item.flags() | _get_qt_item_is_user_checkable())
                        offset_item.setCheckState(0, _get_qt_unchecked())
                        fig_item.addChild(offset_item)
                        offset_added = True
                        
                        if first_implemented_item is None:
                            first_implemented_item = offset_item
                    
                    # If no offsets added, add a placeholder message
                    if not offset_added:
                        if filter_with_spectrum:
                            placeholder = QtWidgets.QTreeWidgetItem(["(No offsets with spectrum loaded)"])
                        else:
                            placeholder = QtWidgets.QTreeWidgetItem(["(No offsets loaded)"])
                        placeholder.setFlags(placeholder.flags() & ~_get_qt_item_is_enabled())
                        placeholder.setForeground(0, QtGui.QBrush(QtGui.QColor(128, 128, 128)))
                        fig_item.addChild(placeholder)
                    
                    # Expand the parent item
                    fig_item.setExpanded(True)
                elif is_implemented:
                    # Regular implemented items get checkboxes
                    fig_item.setFlags(fig_item.flags() | _get_qt_item_is_user_checkable())
                    fig_item.setCheckState(0, _get_qt_unchecked())
                    if first_implemented_item is None:
                        first_implemented_item = fig_item
                else:
                    # Disable unimplemented items
                    fig_item.setFlags(fig_item.flags() & ~_get_qt_item_is_enabled())
                    fig_item.setForeground(0, QtGui.QBrush(QtGui.QColor(128, 128, 128)))

                category_item.addChild(fig_item)

            self.plot_tree.addTopLevelItem(category_item)

        # Check the first implemented item by default and select it
        if first_implemented_item:
            first_implemented_item.setCheckState(0, _get_qt_checked())
            self.plot_tree.setCurrentItem(first_implemented_item)
            self._update_selection_count()

    def _filter_plot_types(self, text: str):
        """Filter the tree view based on search text."""
        search_text = text.lower().strip()

        for i in range(self.plot_tree.topLevelItemCount()):
            category = self.plot_tree.topLevelItem(i)
            category_visible = False

            for j in range(category.childCount()):
                child = category.child(j)
                data = child.data(0, _get_qt_user_role())

                if data:
                    # Search in name and description
                    name_match = search_text in data['name'].lower()
                    desc_match = search_text in data['description'].lower()
                    key_match = search_text in data['key'].lower()

                    visible = not search_text or name_match or desc_match or key_match
                    child.setHidden(not visible)

                    if visible:
                        category_visible = True
                else:
                    child.setHidden(bool(search_text))

            # Hide empty categories, but expand visible ones
            category.setHidden(not category_visible)
            if category_visible and search_text:
                category.setExpanded(True)

    def _on_tree_selection_changed(self):
        """Handle tree selection change to update description panel."""
        selected_items = self.plot_tree.selectedItems()

        if not selected_items:
            self.plot_desc_label.setText("Select a figure type to see its description.")
            self.plot_status_label.setText("")
            self.offset_group.setEnabled(False)
            return

        # Get first selected item's data
        item = selected_items[0]
        data = item.data(0, _get_qt_user_role())

        if data:
            self.plot_desc_label.setText(f"<b>{data['name']}</b><br><br>{data['description']}")

            if data['implemented']:
                self.plot_status_label.setText("<span style='color: green;'>Status: Available</span>")
            else:
                self.plot_status_label.setText("<span style='color: gray;'>Status: Coming Soon</span>")

            # Enable offset options if any checked per-offset type exists
            self._update_offset_options()
        else:
            # Category item selected
            self.plot_desc_label.setText(f"Category: {item.text(0)}")
            self.plot_status_label.setText("")

    def _on_tree_item_changed(self, item, column):
        """Handle tree item checkbox state change."""
        # Update selection count and offset options when checkboxes change
        self._update_selection_count()
        self._update_offset_options()

    def _update_selection_count(self):
        """Update the selection count label."""
        checked_types = self._get_checked_plot_types()
        count = len(checked_types)
        if count == 0:
            self.selection_count_label.setText("No figures selected")
        elif count == 1:
            self.selection_count_label.setText("1 figure selected")
        else:
            self.selection_count_label.setText(f"{count} figures selected")

    def _populate_offset_selector(self):
        """Populate the offset selector combo from loaded layers."""
        self.offset_selector_combo.clear()

        # Get available offsets from controller
        available_offsets = self._get_available_offsets()

        if not available_offsets:
            self.offset_selector_combo.addItem("No offsets loaded", None)
            return

        for offset_info in available_offsets:
            display_text = offset_info['name']
            if offset_info.get('has_spectrum'):
                display_text += " [+spectrum]"
            self.offset_selector_combo.addItem(display_text, offset_info)

    def _populate_grid_offset_list(self):
        """Populate the grid offset selection list with checkboxes."""
        self.grid_offset_list.clear()
        
        available_offsets = self._get_available_offsets()
        
        if not available_offsets:
            item = QtWidgets.QListWidgetItem("No offsets loaded")
            item.setFlags(item.flags() & ~_get_qt_item_is_user_checkable())
            self.grid_offset_list.addItem(item)
            return
        
        for offset_info in available_offsets:
            display_text = offset_info['name']
            if offset_info.get('has_spectrum'):
                display_text += " [+spectrum]"
            
            item = QtWidgets.QListWidgetItem(display_text)
            item.setFlags(item.flags() | _get_qt_item_is_user_checkable())
            item.setCheckState(_get_qt_checked())  # Default: all selected
            item.setData(_get_qt_user_role(), offset_info)
            self.grid_offset_list.addItem(item)
    
    def _select_all_grid_offsets(self):
        """Select all offsets in the grid list."""
        for i in range(self.grid_offset_list.count()):
            item = self.grid_offset_list.item(i)
            if item.flags() & _get_qt_item_is_user_checkable():
                item.setCheckState(_get_qt_checked())
    
    def _deselect_all_grid_offsets(self):
        """Deselect all offsets in the grid list."""
        for i in range(self.grid_offset_list.count()):
            item = self.grid_offset_list.item(i)
            if item.flags() & _get_qt_item_is_user_checkable():
                item.setCheckState(_get_qt_unchecked())
    
    def _get_selected_grid_offsets(self) -> List[int]:
        """Get list of selected offset indices for grid export.
        
        Returns:
            List of offset indices that are checked
        """
        selected_indices = []
        for i in range(self.grid_offset_list.count()):
            item = self.grid_offset_list.item(i)
            if item.checkState() == _get_qt_checked():
                offset_info = item.data(_get_qt_user_role())
                if offset_info and 'index' in offset_info:
                    selected_indices.append(offset_info['index'])
        return selected_indices

    def _get_available_offsets(self) -> List[Dict]:
        """Get list of available offsets from loaded data.
        
        Returns:
            List of dicts: [{'name': 'Offset 12m', 'index': 0, 'has_spectrum': True}, ...]
        """
        offsets = []

        # Get layers from controller if available
        try:
            # Use the correct path: _layers_model instead of model
            if hasattr(self.controller, '_layers_model') and hasattr(self.controller._layers_model, 'layers'):
                layers = self.controller._layers_model.layers
                for i, layer in enumerate(layers):
                    # Use 'label' attribute which contains actual layer name like 'fk +66'
                    layer_name = getattr(layer, 'label', None) or getattr(layer, 'name', None) or f'Offset {i+1}'
                    source_pos = getattr(layer, 'source_position', None)
                    # Check for spectrum_data attribute (the actual spectrum data)
                    has_spectrum = (hasattr(layer, 'spectrum_data') and layer.spectrum_data is not None) or \
                                   (hasattr(layer, 'spectrum') and layer.spectrum is not None)

                    offset_info = {
                        'name': layer_name,
                        'index': i,
                        'source_position': source_pos,
                        'has_spectrum': has_spectrum,
                        'data_points': len(layer.frequency) if hasattr(layer, 'frequency') else 0
                    }
                    offsets.append(offset_info)
        except Exception:
            # If controller not properly initialized, return empty list
            pass

        return offsets

    def _update_offset_options(self):
        """Enable/disable offset options based on checked items."""
        checked_types = self._get_checked_plot_types()

        # Check if any offset-related types are checked
        # Note: Individual offset types (curve_only, with_spectrum, spectrum_only)
        # now have sub-items so we check for per_offset and offset_grid
        offset_keys = ['per_offset', 'offset_grid']
        needs_offset = any(
            any(offset_key in key for offset_key in offset_keys) 
            for key, _, _ in checked_types
        )

        self.offset_group.setEnabled(needs_offset)

        # Check if offset_grid is selected
        is_grid_selected = any('offset_grid' in key for key, _, _ in checked_types)
        
        # Show/hide appropriate UI elements based on figure type
        # Grid selection UI (for Comparison Grid)
        self.grid_offset_selection_label.setVisible(is_grid_selected)
        self.select_all_btn.setVisible(is_grid_selected)
        self.deselect_all_btn.setVisible(is_grid_selected)
        self.grid_offset_list.setVisible(is_grid_selected)
        
        # Single offset selector (for individual offset types)
        is_single_offset = any(
            key in ['offset_curve_only', 'offset_with_spectrum', 'offset_spectrum_only']
            for key, _, _ in checked_types
        )
        self.offset_selector_label.setVisible(is_single_offset)
        self.offset_selector_combo.setVisible(is_single_offset)

        # Refresh offset selector when group becomes enabled
        if needs_offset:
            self._populate_offset_selector()
            self._populate_grid_offset_list()
        
    def _get_checked_plot_types(self) -> List[Tuple[str, str, Optional[Dict]]]:
        """Get list of checked plot types from tree view (using checkboxes).
        
        Returns:
            List of tuples (internal_key, suffix_for_filename, offset_info_or_None)
            For offset types, offset_info contains 'offset_index', 'offset_name', etc.
        """
        checked_types = []

        for i in range(self.plot_tree.topLevelItemCount()):
            category = self.plot_tree.topLevelItem(i)
            for j in range(category.childCount()):
                child = category.child(j)
                data = child.data(0, _get_qt_user_role())
                
                if data and data.get('implemented', False):
                    # Check if item is checked (for regular items)
                    if child.checkState(0) == _get_qt_checked():
                        key = data['key']
                        checked_types.append((key, key, None))
                
                # Check grandchildren (offset sub-items)
                for k in range(child.childCount()):
                    grandchild = child.child(k)
                    grandchild_data = grandchild.data(0, _get_qt_user_role())
                    
                    if grandchild_data and grandchild_data.get('implemented', False):
                        if grandchild.checkState(0) == _get_qt_checked():
                            key = grandchild_data['key']
                            offset_name = grandchild_data.get('offset_name', '')
                            # Create unique suffix for filename
                            suffix = f"{key}_{offset_name.replace(' ', '_').lower()}"
                            
                            offset_info = {
                                'offset_index': grandchild_data.get('offset_index'),
                                'offset_name': grandchild_data.get('offset_name'),
                                'has_spectrum': grandchild_data.get('has_spectrum', False)
                            }
                            checked_types.append((key, suffix, offset_info))

        return checked_types

    def _get_selected_plot_types(self) -> List[Tuple[str, str, Optional[Dict]]]:
        """Get list of selected plot types from tree view.
        
        Now uses checkbox state instead of selection.
        
        Returns:
            List of tuples (internal_key, suffix_for_filename, offset_info_or_None)
        """
        return self._get_checked_plot_types()

    def _build_styling_tab(self):
        """Build the Styling options tab."""
        tab = QtWidgets.QWidget()
        layout = QtWidgets.QVBoxLayout(tab)

        # Use scroll area for styling tab
        scroll = QtWidgets.QScrollArea()
        scroll.setWidgetResizable(True)
        scroll_content = QtWidgets.QWidget()
        scroll_layout = QtWidgets.QVBoxLayout(scroll_content)

        # Journal Preset (NEW)
        preset_group = QtWidgets.QGroupBox("Journal Preset")
        preset_layout = QtWidgets.QFormLayout(preset_group)

        self.journal_preset_combo = QtWidgets.QComboBox()
        self.journal_preset_combo.addItems(['Custom', 'Nature', 'AGU', 'Geophysics', 'Presentation', 'Poster'])
        self.journal_preset_combo.currentTextChanged.connect(self._on_preset_changed)
        preset_layout.addRow("Preset:", self.journal_preset_combo)

        preset_note = QtWidgets.QLabel("Presets auto-configure figure size, DPI, and fonts")
        preset_note.setStyleSheet("color: gray; font-style: italic;")
        preset_layout.addRow("", preset_note)

        scroll_layout.addWidget(preset_group)

        # Title Settings (NEW)
        title_group = QtWidgets.QGroupBox("Title")
        title_layout = QtWidgets.QFormLayout(title_group)

        self.title_edit = QtWidgets.QLineEdit()
        self.title_edit.setPlaceholderText("Optional figure title")
        title_layout.addRow("Title:", self.title_edit)

        self.title_fontsize_spin = QtWidgets.QSpinBox()
        self.title_fontsize_spin.setRange(8, 24)
        self.title_fontsize_spin.setValue(14)
        title_layout.addRow("Title font size:", self.title_fontsize_spin)

        scroll_layout.addWidget(title_group)

        # Figure settings
        fig_group = QtWidgets.QGroupBox("Figure Settings")
        fig_layout = QtWidgets.QFormLayout(fig_group)

        self.figsize_width_spin = QtWidgets.QDoubleSpinBox()
        self.figsize_width_spin.setRange(2.0, 20.0)
        self.figsize_width_spin.setValue(8.0)
        self.figsize_width_spin.setSingleStep(0.5)
        fig_layout.addRow("Width (inches):", self.figsize_width_spin)

        self.figsize_height_spin = QtWidgets.QDoubleSpinBox()
        self.figsize_height_spin.setRange(2.0, 20.0)
        self.figsize_height_spin.setValue(6.0)
        self.figsize_height_spin.setSingleStep(0.5)
        fig_layout.addRow("Height (inches):", self.figsize_height_spin)

        self.dpi_spin = QtWidgets.QSpinBox()
        self.dpi_spin.setRange(72, 2400)
        self.dpi_spin.setValue(300)
        self.dpi_spin.setSingleStep(50)
        fig_layout.addRow("DPI:", self.dpi_spin)

        scroll_layout.addWidget(fig_group)

        # Font settings
        font_group = QtWidgets.QGroupBox("Font Settings")
        font_layout = QtWidgets.QFormLayout(font_group)

        self.font_family_combo = QtWidgets.QComboBox()
        self.font_family_combo.addItems(['serif', 'sans-serif', 'monospace', 'Times New Roman', 'Arial', 'Helvetica'])
        font_layout.addRow("Font family:", self.font_family_combo)

        self.font_size_spin = QtWidgets.QSpinBox()
        self.font_size_spin.setRange(6, 24)
        self.font_size_spin.setValue(11)
        font_layout.addRow("Font size:", self.font_size_spin)

        self.font_bold_check = QtWidgets.QCheckBox("Bold")
        self.font_bold_check.setChecked(False)
        font_layout.addRow("Font weight:", self.font_bold_check)

        scroll_layout.addWidget(font_group)

        # Line/marker settings
        line_group = QtWidgets.QGroupBox("Line & Marker Settings")
        line_layout = QtWidgets.QFormLayout(line_group)

        self.line_width_spin = QtWidgets.QDoubleSpinBox()
        self.line_width_spin.setRange(0.5, 5.0)
        self.line_width_spin.setValue(1.5)
        self.line_width_spin.setSingleStep(0.1)
        line_layout.addRow("Line width:", self.line_width_spin)

        self.marker_size_spin = QtWidgets.QDoubleSpinBox()
        self.marker_size_spin.setRange(1.0, 10.0)
        self.marker_size_spin.setValue(4.0)
        self.marker_size_spin.setSingleStep(0.5)
        line_layout.addRow("Marker size:", self.marker_size_spin)

        # Marker style (NEW)
        self.marker_style_combo = QtWidgets.QComboBox()
        self.marker_style_combo.addItems(['o (circle)', 's (square)', '^ (triangle)', 'D (diamond)', 'x (cross)', '+ (plus)', '. (point)'])
        line_layout.addRow("Marker style:", self.marker_style_combo)

        scroll_layout.addWidget(line_group)

        # Legend settings (NEW)
        legend_group = QtWidgets.QGroupBox("Legend")
        legend_layout = QtWidgets.QFormLayout(legend_group)

        self.legend_position_combo = QtWidgets.QComboBox()
        self.legend_position_combo.addItems(['best', 'upper left', 'upper right', 'lower left', 'lower right', 
                                              'center left', 'center right', 'upper center', 'lower center',
                                              'outside right', 'outside top', 'outside bottom'])
        legend_layout.addRow("Position:", self.legend_position_combo)

        self.legend_columns_spin = QtWidgets.QSpinBox()
        self.legend_columns_spin.setRange(1, 5)
        self.legend_columns_spin.setValue(1)
        legend_layout.addRow("Columns:", self.legend_columns_spin)

        self.legend_frameon_check = QtWidgets.QCheckBox("Show legend frame")
        self.legend_frameon_check.setChecked(False)
        legend_layout.addRow("", self.legend_frameon_check)

        scroll_layout.addWidget(legend_group)

        # Color settings
        color_group = QtWidgets.QGroupBox("Color Settings")
        color_layout = QtWidgets.QFormLayout(color_group)

        self.color_palette_combo = QtWidgets.QComboBox()
        self.color_palette_combo.addItems(['vibrant', 'muted', 'bright', 'high_contrast'])
        color_layout.addRow("Color palette:", self.color_palette_combo)

        palette_note = QtWidgets.QLabel("All palettes are colorblind-friendly")
        palette_note.setStyleSheet("color: gray; font-style: italic;")
        color_layout.addRow("", palette_note)

        self.uncertainty_alpha_spin = QtWidgets.QDoubleSpinBox()
        self.uncertainty_alpha_spin.setRange(0.0, 1.0)
        self.uncertainty_alpha_spin.setValue(0.3)
        self.uncertainty_alpha_spin.setSingleStep(0.05)
        color_layout.addRow("Uncertainty alpha:", self.uncertainty_alpha_spin)

        scroll_layout.addWidget(color_group)

        # Spectrum Options (for Source Offset Analysis)
        spectrum_group = QtWidgets.QGroupBox("Spectrum Options")
        spectrum_layout = QtWidgets.QFormLayout(spectrum_group)

        self.spectrum_colormap_combo = QtWidgets.QComboBox()
        self.spectrum_colormap_combo.addItems([
            'viridis', 'plasma', 'inferno', 'magma', 'cividis',
            'jet', 'hot', 'coolwarm', 'RdYlBu', 'Spectral', 'turbo'
        ])
        spectrum_layout.addRow("Colormap:", self.spectrum_colormap_combo)

        self.spectrum_render_mode_combo = QtWidgets.QComboBox()
        self.spectrum_render_mode_combo.addItems(['imshow (fast)', 'contour (smooth)'])
        spectrum_layout.addRow("Render mode:", self.spectrum_render_mode_combo)

        self.spectrum_alpha_spin = QtWidgets.QDoubleSpinBox()
        self.spectrum_alpha_spin.setRange(0.0, 1.0)
        self.spectrum_alpha_spin.setValue(0.8)
        self.spectrum_alpha_spin.setSingleStep(0.1)
        spectrum_layout.addRow("Transparency:", self.spectrum_alpha_spin)

        self.spectrum_levels_spin = QtWidgets.QSpinBox()
        self.spectrum_levels_spin.setRange(10, 100)
        self.spectrum_levels_spin.setValue(30)
        spectrum_layout.addRow("Contour levels:", self.spectrum_levels_spin)

        # Colorbar orientation option (replaces simple show/hide checkbox)
        self.colorbar_orientation_combo = QtWidgets.QComboBox()
        self.colorbar_orientation_combo.addItems(['Vertical (Right)', 'Horizontal (Bottom)', 'None (Hidden)'])
        spectrum_layout.addRow("Colorbar:", self.colorbar_orientation_combo)

        spectrum_note = QtWidgets.QLabel("Used for offset analysis with spectrum background")
        spectrum_note.setStyleSheet("color: gray; font-style: italic;")
        spectrum_layout.addRow("", spectrum_note)

        scroll_layout.addWidget(spectrum_group)

        # Peak Overlay Options (for spectrum plots)
        peak_group = QtWidgets.QGroupBox("Peak/Curve Overlay")
        peak_layout = QtWidgets.QFormLayout(peak_group)

        self.peak_color_combo = QtWidgets.QComboBox()
        self.peak_color_combo.addItems(['white', 'black', 'yellow', 'cyan', 'magenta', 'lime', 'red'])
        peak_layout.addRow("Curve color:", self.peak_color_combo)

        self.peak_outline_check = QtWidgets.QCheckBox("Show contrast outline")
        self.peak_outline_check.setChecked(True)
        peak_layout.addRow("", self.peak_outline_check)

        self.peak_line_width_spin = QtWidgets.QDoubleSpinBox()
        self.peak_line_width_spin.setRange(1.0, 6.0)
        self.peak_line_width_spin.setValue(2.5)
        self.peak_line_width_spin.setSingleStep(0.5)
        peak_layout.addRow("Line width:", self.peak_line_width_spin)

        # Curve overlay style (line/markers/both)
        self.curve_overlay_style_combo = QtWidgets.QComboBox()
        self.curve_overlay_style_combo.addItems(['Line Only', 'Markers Only', 'Line + Markers'])
        peak_layout.addRow("Overlay style:", self.curve_overlay_style_combo)

        peak_note = QtWidgets.QLabel("Styling for curves overlaid on spectrum background")
        peak_note.setStyleSheet("color: gray; font-style: italic;")
        peak_layout.addRow("", peak_note)

        scroll_layout.addWidget(peak_group)

        scroll_layout.addStretch()
        scroll.setWidget(scroll_content)
        layout.addWidget(scroll)
        self.tabs.addTab(tab, "Styling")

    def _on_preset_changed(self, preset_name: str):
        """Apply journal preset settings."""
        presets = {
            'Nature': {'width': 3.5, 'height': 3.5, 'dpi': 300, 'font': 'sans-serif', 'font_size': 8},
            'AGU': {'width': 3.74, 'height': 4.53, 'dpi': 300, 'font': 'sans-serif', 'font_size': 9},
            'Geophysics': {'width': 3.5, 'height': 4.0, 'dpi': 600, 'font': 'sans-serif', 'font_size': 9},
            'Presentation': {'width': 10, 'height': 7.5, 'dpi': 150, 'font': 'sans-serif', 'font_size': 14},
            'Poster': {'width': 12, 'height': 9, 'dpi': 150, 'font': 'sans-serif', 'font_size': 16},
        }
        if preset_name in presets:
            p = presets[preset_name]
            self.figsize_width_spin.setValue(p['width'])
            self.figsize_height_spin.setValue(p['height'])
            self.dpi_spin.setValue(p['dpi'])
            idx = self.font_family_combo.findText(p['font'])
            if idx >= 0:
                self.font_family_combo.setCurrentIndex(idx)
            self.font_size_spin.setValue(p['font_size'])

    def _build_axes_tab(self):
        """Build the Axes configuration tab."""
        tab = QtWidgets.QWidget()
        layout = QtWidgets.QVBoxLayout(tab)

        # Grid settings
        grid_group = QtWidgets.QGroupBox("Grid")
        grid_layout = QtWidgets.QVBoxLayout(grid_group)

        self.show_grid_check = QtWidgets.QCheckBox("Show grid")
        self.show_grid_check.setChecked(True)
        grid_layout.addWidget(self.show_grid_check)

        grid_alpha_layout = QtWidgets.QHBoxLayout()
        grid_alpha_layout.addWidget(QtWidgets.QLabel("Grid alpha:"))
        self.grid_alpha_spin = QtWidgets.QDoubleSpinBox()
        self.grid_alpha_spin.setRange(0.0, 1.0)
        self.grid_alpha_spin.setValue(0.3)
        self.grid_alpha_spin.setSingleStep(0.05)
        grid_alpha_layout.addWidget(self.grid_alpha_spin)
        grid_alpha_layout.addStretch()
        grid_layout.addLayout(grid_alpha_layout)

        layout.addWidget(grid_group)

        # Labels
        labels_group = QtWidgets.QGroupBox("Axis Labels")
        labels_layout = QtWidgets.QFormLayout(labels_group)

        self.xlabel_edit = QtWidgets.QLineEdit("Frequency (Hz)")
        labels_layout.addRow("X-axis label:", self.xlabel_edit)

        self.ylabel_edit = QtWidgets.QLineEdit("Phase Velocity (m/s)")
        labels_layout.addRow("Y-axis label:", self.ylabel_edit)

        layout.addWidget(labels_group)

        # Limits
        limits_group = QtWidgets.QGroupBox("Axis Limits")
        limits_layout = QtWidgets.QFormLayout(limits_group)

        self.xlim_auto_check = QtWidgets.QCheckBox("Auto")
        self.xlim_auto_check.setChecked(True)
        limits_layout.addRow("X-axis:", self.xlim_auto_check)

        xlim_layout = QtWidgets.QHBoxLayout()
        self.xlim_min_spin = QtWidgets.QDoubleSpinBox()
        self.xlim_min_spin.setRange(0.1, 1000.0)
        self.xlim_min_spin.setValue(1.0)
        self.xlim_min_spin.setEnabled(False)
        xlim_layout.addWidget(QtWidgets.QLabel("Min:"))
        xlim_layout.addWidget(self.xlim_min_spin)
        self.xlim_max_spin = QtWidgets.QDoubleSpinBox()
        self.xlim_max_spin.setRange(0.1, 1000.0)
        self.xlim_max_spin.setValue(100.0)
        self.xlim_max_spin.setEnabled(False)
        xlim_layout.addWidget(QtWidgets.QLabel("Max:"))
        xlim_layout.addWidget(self.xlim_max_spin)
        xlim_layout.addStretch()
        limits_layout.addRow("", xlim_layout)

        self.ylim_auto_check = QtWidgets.QCheckBox("Auto")
        self.ylim_auto_check.setChecked(True)
        limits_layout.addRow("Y-axis:", self.ylim_auto_check)

        ylim_layout = QtWidgets.QHBoxLayout()
        self.ylim_min_spin = QtWidgets.QDoubleSpinBox()
        self.ylim_min_spin.setRange(0.0, 10000.0)
        self.ylim_min_spin.setValue(0.0)
        self.ylim_min_spin.setEnabled(False)
        ylim_layout.addWidget(QtWidgets.QLabel("Min:"))
        ylim_layout.addWidget(self.ylim_min_spin)
        self.ylim_max_spin = QtWidgets.QDoubleSpinBox()
        self.ylim_max_spin.setRange(0.0, 10000.0)
        self.ylim_max_spin.setValue(1000.0)
        self.ylim_max_spin.setEnabled(False)
        ylim_layout.addWidget(QtWidgets.QLabel("Max:"))
        ylim_layout.addWidget(self.ylim_max_spin)
        ylim_layout.addStretch()
        limits_layout.addRow("", ylim_layout)

        layout.addWidget(limits_group)

        # Connect auto checkboxes
        self.xlim_auto_check.toggled.connect(lambda checked: self.xlim_min_spin.setDisabled(checked))
        self.xlim_auto_check.toggled.connect(lambda checked: self.xlim_max_spin.setDisabled(checked))
        self.ylim_auto_check.toggled.connect(lambda checked: self.ylim_min_spin.setDisabled(checked))
        self.ylim_auto_check.toggled.connect(lambda checked: self.ylim_max_spin.setDisabled(checked))

        # Near-field marking
        nf_group = QtWidgets.QGroupBox("Near-Field Marking")
        nf_layout = QtWidgets.QVBoxLayout(nf_group)

        self.mark_near_field_check = QtWidgets.QCheckBox("Mark near-field data")
        self.mark_near_field_check.setChecked(True)
        nf_layout.addWidget(self.mark_near_field_check)

        nf_style_layout = QtWidgets.QHBoxLayout()
        nf_style_layout.addWidget(QtWidgets.QLabel("Style:"))
        self.nf_style_combo = QtWidgets.QComboBox()
        self.nf_style_combo.addItems(['faded', 'crossed', 'none'])
        nf_style_layout.addWidget(self.nf_style_combo)
        nf_style_layout.addStretch()
        nf_layout.addLayout(nf_style_layout)

        nf_thresh_layout = QtWidgets.QHBoxLayout()
        nf_thresh_layout.addWidget(QtWidgets.QLabel("NACD threshold:"))
        self.nacd_thresh_spin = QtWidgets.QDoubleSpinBox()
        self.nacd_thresh_spin.setRange(0.1, 10.0)
        self.nacd_thresh_spin.setValue(1.0)
        self.nacd_thresh_spin.setSingleStep(0.1)
        nf_thresh_layout.addWidget(self.nacd_thresh_spin)
        nf_thresh_layout.addStretch()
        nf_layout.addLayout(nf_thresh_layout)

        nf_alpha_layout = QtWidgets.QHBoxLayout()
        nf_alpha_layout.addWidget(QtWidgets.QLabel("Near-field alpha:"))
        self.nf_alpha_spin = QtWidgets.QDoubleSpinBox()
        self.nf_alpha_spin.setRange(0.0, 1.0)
        self.nf_alpha_spin.setValue(0.4)
        self.nf_alpha_spin.setSingleStep(0.05)
        nf_alpha_layout.addWidget(self.nf_alpha_spin)
        nf_alpha_layout.addStretch()
        nf_layout.addLayout(nf_alpha_layout)

        layout.addWidget(nf_group)

        layout.addStretch()
        self.tabs.addTab(tab, "Axes & Limits")

    def _build_output_tab(self):
        """Build the Output settings tab."""
        tab = QtWidgets.QWidget()
        layout = QtWidgets.QVBoxLayout(tab)

        # Output directory
        dir_group = QtWidgets.QGroupBox("Output Directory")
        dir_layout = QtWidgets.QVBoxLayout(dir_group)

        dir_select_layout = QtWidgets.QHBoxLayout()
        self.output_dir_edit = QtWidgets.QLineEdit()
        self.output_dir_edit.setPlaceholderText("Select output directory...")
        dir_select_layout.addWidget(self.output_dir_edit)

        self.browse_button = QtWidgets.QPushButton("Browse...")
        self.browse_button.clicked.connect(self._on_browse)
        dir_select_layout.addWidget(self.browse_button)

        dir_layout.addLayout(dir_select_layout)

        # Project name (used as filename prefix)
        filename_layout = QtWidgets.QHBoxLayout()
        filename_layout.addWidget(QtWidgets.QLabel("Project Name:"))
        self.filename_edit = QtWidgets.QLineEdit("dispersion_curve")
        self.filename_edit.setPlaceholderText("Enter project name (used as filename prefix)")
        filename_layout.addWidget(self.filename_edit)
        dir_layout.addLayout(filename_layout)

        note_label = QtWidgets.QLabel(
            "Note: When multiple figures are selected, files are named as:\n"
            "<project_name>_<figure_type>.<extension>\n"
            "e.g., dispersion_curve_aggregated.pdf, dispersion_curve_per_offset.pdf"
        )
        note_label.setStyleSheet("color: gray; font-style: italic; font-size: 9pt;")
        note_label.setWordWrap(True)
        dir_layout.addWidget(note_label)

        layout.addWidget(dir_group)

        # Format
        format_group = QtWidgets.QGroupBox("Output Format")
        format_layout = QtWidgets.QVBoxLayout(format_group)

        self.format_pdf = QtWidgets.QRadioButton("PDF (vector, recommended)")
        self.format_png = QtWidgets.QRadioButton("PNG (raster)")
        self.format_svg = QtWidgets.QRadioButton("SVG (vector)")
        self.format_eps = QtWidgets.QRadioButton("EPS (vector)")
        self.format_pptx = QtWidgets.QRadioButton("PPTX (PowerPoint)")

        self.format_pdf.setChecked(True)

        format_layout.addWidget(self.format_pdf)
        format_layout.addWidget(self.format_png)
        format_layout.addWidget(self.format_svg)
        format_layout.addWidget(self.format_eps)
        format_layout.addWidget(self.format_pptx)
        
        # PPTX-specific options (only visible when PPTX selected)
        self.pptx_options_widget = QtWidgets.QWidget()
        pptx_options_layout = QtWidgets.QVBoxLayout(self.pptx_options_widget)
        pptx_options_layout.setContentsMargins(20, 0, 0, 0)  # Indent
        
        self.pptx_combine_check = QtWidgets.QCheckBox("Combine all figures into one PPTX file")
        self.pptx_combine_check.setChecked(True)
        pptx_options_layout.addWidget(self.pptx_combine_check)
        
        self.pptx_grid_slides_check = QtWidgets.QCheckBox("For grids: create individual slides + combined slide")
        self.pptx_grid_slides_check.setChecked(True)
        pptx_options_layout.addWidget(self.pptx_grid_slides_check)
        
        self.pptx_options_widget.setVisible(False)
        format_layout.addWidget(self.pptx_options_widget)
        
        # Connect format radio buttons to show/hide PPTX options
        self.format_pptx.toggled.connect(self.pptx_options_widget.setVisible)

        layout.addWidget(format_group)

        # Layout options
        layout_group = QtWidgets.QGroupBox("Layout")
        layout_layout = QtWidgets.QVBoxLayout(layout_group)

        self.tight_layout_check = QtWidgets.QCheckBox("Use tight layout")
        self.tight_layout_check.setChecked(True)
        layout_layout.addWidget(self.tight_layout_check)

        layout.addWidget(layout_group)

        layout.addStretch()
        self.tabs.addTab(tab, "Output")

    def _on_browse(self):
        """Open dialog to select output directory."""
        # Get current directory if set
        current_dir = self.output_dir_edit.text().strip()
        if not current_dir:
            current_dir = str(Path.home())

        # Handle Qt version compatibility for ShowDirsOnly option
        try:
            # Try newer Qt6 enum style first
            options = QtWidgets.QFileDialog.Option.ShowDirsOnly
        except AttributeError:
            try:
                # Fall back to Qt5 style
                options = QtWidgets.QFileDialog.ShowDirsOnly
            except AttributeError:
                # If neither works, use no options
                options = None

        try:
            if options is not None:
                dir_path = QtWidgets.QFileDialog.getExistingDirectory(
                    self,
                    "Select Output Directory",
                    current_dir,
                    options
                )
            else:
                dir_path = QtWidgets.QFileDialog.getExistingDirectory(
                    self,
                    "Select Output Directory",
                    current_dir
                )
        except Exception as e:
            # Fallback without options if there's any error
            dir_path = QtWidgets.QFileDialog.getExistingDirectory(
                self,
                "Select Output Directory",
                current_dir
            )

        if dir_path:
            self.output_dir_edit.setText(dir_path)

    def _gather_config(self) -> PlotConfig:
        """Gather configuration from UI widgets."""
        # Determine format
        if self.format_pdf.isChecked():
            output_format = 'pdf'
        elif self.format_png.isChecked():
            output_format = 'png'
        elif self.format_svg.isChecked():
            output_format = 'svg'
        elif self.format_pptx.isChecked():
            output_format = 'pptx'
        else:
            output_format = 'eps'

        # Determine limits
        xlim = None
        if not self.xlim_auto_check.isChecked():
            xlim = (self.xlim_min_spin.value(), self.xlim_max_spin.value())

        ylim = None
        if not self.ylim_auto_check.isChecked():
            ylim = (self.ylim_min_spin.value(), self.ylim_max_spin.value())

        # Get marker style (extract first character)
        marker_text = self.marker_style_combo.currentText()
        marker_style = marker_text.split(' ')[0] if marker_text else 'o'

        # Get title (empty string = None)
        title = self.title_edit.text().strip() or None

        config = PlotConfig(
            figsize=(self.figsize_width_spin.value(), self.figsize_height_spin.value()),
            dpi=self.dpi_spin.value(),
            font_family=self.font_family_combo.currentText(),
            font_size=self.font_size_spin.value(),
            font_weight='bold' if self.font_bold_check.isChecked() else 'normal',
            line_width=self.line_width_spin.value(),
            marker_size=self.marker_size_spin.value(),
            marker_style=marker_style,
            title=title,
            title_fontsize=self.title_fontsize_spin.value(),
            legend_position=self.legend_position_combo.currentText(),
            legend_columns=self.legend_columns_spin.value(),
            legend_frameon=self.legend_frameon_check.isChecked(),
            color_palette=self.color_palette_combo.currentText(),
            uncertainty_alpha=self.uncertainty_alpha_spin.value(),
            near_field_alpha=self.nf_alpha_spin.value(),
            mark_near_field=self.mark_near_field_check.isChecked(),
            near_field_style=self.nf_style_combo.currentText(),
            nacd_threshold=self.nacd_thresh_spin.value(),
            show_grid=self.show_grid_check.isChecked(),
            grid_alpha=self.grid_alpha_spin.value(),
            xlabel=self.xlabel_edit.text(),
            ylabel=self.ylabel_edit.text(),
            xlim=xlim,
            ylim=ylim,
            output_format=output_format,
            tight_layout=self.tight_layout_check.isChecked(),
            # Spectrum options
            spectrum_colormap=self.spectrum_colormap_combo.currentText(),
            spectrum_render_mode='contour' if 'contour' in self.spectrum_render_mode_combo.currentText().lower() else 'imshow',
            spectrum_alpha=self.spectrum_alpha_spin.value(),
            spectrum_levels=self.spectrum_levels_spin.value(),
            show_spectrum_colorbar=self._get_colorbar_orientation() != 'none',
            spectrum_colorbar_orientation=self._get_colorbar_orientation(),
            # Peak overlay options
            peak_color=self.peak_color_combo.currentText(),
            peak_outline=self.peak_outline_check.isChecked(),
            peak_line_width=self.peak_line_width_spin.value(),
            # Curve overlay style
            curve_overlay_style=self._get_curve_overlay_style(),
            # Grid options
            grid_offset_indices=self._get_selected_grid_offsets() or None,
        )

        return config

    def _get_colorbar_orientation(self) -> str:
        """Get colorbar orientation from combo box."""
        text = self.colorbar_orientation_combo.currentText()
        if 'Vertical' in text:
            return 'vertical'
        elif 'Horizontal' in text:
            return 'horizontal'
        return 'none'

    def _get_curve_overlay_style(self) -> str:
        """Convert combo text to curve overlay style value."""
        text = self.curve_overlay_style_combo.currentText()
        if 'Markers Only' in text:
            return 'markers'
        elif 'Line + Markers' in text:
            return 'line+markers'
        else:
            return 'line'

    def _on_generate(self):
        """Generate the publication figure(s)."""
        # Validate output directory
        output_dir = self.output_dir_edit.text().strip()
        if not output_dir:
            QtWidgets.QMessageBox.warning(
                self,
                "No Output Directory",
                "Please select an output directory."
            )
            return

        # Validate directory exists
        dir_path = Path(output_dir)
        if not dir_path.exists():
            QtWidgets.QMessageBox.warning(
                self,
                "Directory Not Found",
                f"The directory does not exist:\n{output_dir}"
            )
            return

        if not dir_path.is_dir():
            QtWidgets.QMessageBox.warning(
                self,
                "Invalid Directory",
                f"The path is not a directory:\n{output_dir}"
            )
            return

        # Get filename
        filename = self.filename_edit.text().strip()
        if not filename:
            QtWidgets.QMessageBox.warning(
                self,
                "No Filename",
                "Please enter a filename."
            )
            return

        # Determine file extension based on format
        if self.format_pdf.isChecked():
            extension = ".pdf"
        elif self.format_png.isChecked():
            extension = ".png"
        elif self.format_svg.isChecked():
            extension = ".svg"
        elif self.format_pptx.isChecked():
            extension = ".pptx"
        else:
            extension = ".eps"

        # Check if PPTX is selected
        is_pptx = self.format_pptx.isChecked()
        pptx_combine = is_pptx and self.pptx_combine_check.isChecked()
        pptx_grid_slides = is_pptx and self.pptx_grid_slides_check.isChecked()

        # Remove extension from filename if user added one
        filename_base = Path(filename).stem

        # Gather configuration
        config = self._gather_config()

        # Create generator from controller
        try:
            generator = PublicationFigureGenerator.from_controller(self.controller)
        except Exception as e:
            QtWidgets.QMessageBox.critical(
                self,
                "Error",
                f"Failed to create figure generator:\n{str(e)}"
            )
            return

        # Get selected plot types from tree view
        selected_types = self._get_selected_plot_types()

        # Validate at least one type is selected
        if not selected_types:
            QtWidgets.QMessageBox.warning(
                self,
                "No Plot Type Selected",
                "Please select at least one plot type to generate."
            )
            return

        # Build file list - now handles (plot_type, suffix, offset_info) tuples
        files_to_create = []
        for plot_type, suffix, offset_info in selected_types:
            # If only one type selected, don't add suffix
            if len(selected_types) == 1:
                output_path = str(dir_path / f"{filename_base}{extension}")
            else:
                output_path = str(dir_path / f"{filename_base}_{suffix}{extension}")
            files_to_create.append((plot_type, output_path, offset_info))

        # Check for existing files
        if pptx_combine:
            # For combined PPTX, only one output file
            combined_pptx_path = str(dir_path / f"{filename_base}{extension}")
            if Path(combined_pptx_path).exists():
                reply = QtWidgets.QMessageBox.question(
                    self,
                    "File Exists",
                    f"The file already exists:\n{combined_pptx_path}\n\nOverwrite?",
                    _get_qt_msgbox_yes() | _get_qt_msgbox_no(),
                    _get_qt_msgbox_no()
                )
                if reply == _get_qt_msgbox_no():
                    return
        else:
            existing_files = [p for _, p, _ in files_to_create if Path(p).exists()]
            if existing_files:
                file_list = '\n'.join(existing_files)
                reply = QtWidgets.QMessageBox.question(
                    self,
                    "Files Exist",
                    f"The following files already exist:\n{file_list}\n\nOverwrite?",
                    _get_qt_msgbox_yes() | _get_qt_msgbox_no(),
                    _get_qt_msgbox_no()
                )
                if reply == _get_qt_msgbox_no():
                    return

        # Generate all selected plots
        try:
            if is_pptx and pptx_combine:
                # Generate all figures to a single PPTX
                generated_files = self._generate_combined_pptx(
                    generator, files_to_create, combined_pptx_path, config
                )
            elif is_pptx:
                # Generate separate PPTX files for each figure
                generated_files = []
                for plot_type, output_path, offset_info in files_to_create:
                    self._generate_single_pptx(generator, plot_type, output_path, config, offset_info)
                    generated_files.append(output_path)
            else:
                # Regular image format export
                generated_files = []
                for plot_type, output_path, offset_info in files_to_create:
                    self._generate_single_plot(generator, plot_type, output_path, config, offset_info)
                    generated_files.append(output_path)

            # Success message
            file_list = '\n'.join(generated_files)
            if len(generated_files) == 1:
                msg = f"Publication figure saved to:\n{file_list}"
            else:
                msg = f"Publication figures saved:\n{file_list}"
            QtWidgets.QMessageBox.information(
                self,
                "Success",
                msg
            )

            # Don't close dialog - user may want to generate more figures
            # self.accept()

        except Exception as e:
            QtWidgets.QMessageBox.critical(
                self,
                "Error",
                f"Failed to generate figure:\n{str(e)}"
            )
    
    def _generate_combined_pptx(self, generator, files_to_create, output_path: str, config: PlotConfig) -> List[str]:
        """Generate all figures into a single PowerPoint file.
        
        Args:
            generator: PublicationFigureGenerator instance
            files_to_create: List of (plot_type, _, offset_info) tuples
            output_path: Path to save the combined PPTX
            config: PlotConfig instance
            
        Returns:
            List containing the output path
        """
        import tempfile
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint export.\n"
                "Install it with: pip install python-pptx"
            )
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio (standard widescreen)
        prs.slide_height = Inches(7.5)
        
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        
        # Build a lookup for display names from FIGURE_TYPES
        display_names = {}
        for category_types in FIGURE_TYPES.values():
            for display_name, internal_key, _, _ in category_types:
                display_names[internal_key] = display_name
        
        with tempfile.TemporaryDirectory() as temp_dir:
            slide_count = 0
            for plot_type, _, offset_info in files_to_create:
                # Generate figure to temp PNG
                temp_png = Path(temp_dir) / f"{plot_type}_{slide_count}.png"
                
                # Temporarily override format to PNG
                original_format = config.output_format
                temp_config = PlotConfig(**{**config.__dict__, 'output_format': 'png'})
                
                self._generate_single_plot(generator, plot_type, str(temp_png), temp_config, offset_info)
                
                # Add slide with image and title
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Create title text
                base_name = display_names.get(plot_type, plot_type.replace('_', ' ').title())
                if offset_info and offset_info.get('offset_name'):
                    slide_title = f"{base_name}: {offset_info['offset_name']}"
                else:
                    slide_title = base_name
                
                # Add title text box
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.5)
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = slide_title
                title_para.font.name = "Times New Roman"
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER
                
                # Add image below title
                left = Inches(0.5)
                top = Inches(0.8)
                width = Inches(12.333)
                
                slide.shapes.add_picture(str(temp_png), left, top, width=width)
                slide_count += 1
        
        prs.save(output_path)
        return [output_path]
    
    def _generate_single_pptx(self, generator, plot_type: str, output_path: str, config: PlotConfig,
                              offset_info: Optional[Dict] = None):
        """Generate a single figure as a PowerPoint file with one slide.
        
        Args:
            generator: PublicationFigureGenerator instance
            plot_type: Internal key of the plot type
            output_path: Path to save the PPTX
            config: PlotConfig instance
            offset_info: Optional dict with offset details
        """
        import tempfile
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint export.\n"
                "Install it with: pip install python-pptx"
            )
        
        # Build a lookup for display names from FIGURE_TYPES
        display_names = {}
        for category_types in FIGURE_TYPES.values():
            for display_name, internal_key, _, _ in category_types:
                display_names[internal_key] = display_name
        
        # Create presentation with single slide
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        blank_slide_layout = prs.slide_layouts[6]
        
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_png = Path(temp_dir) / "figure.png"
            
            # Generate figure to temp PNG
            temp_config = PlotConfig(**{**config.__dict__, 'output_format': 'png'})
            self._generate_single_plot(generator, plot_type, str(temp_png), temp_config, offset_info)
            
            # Add slide with image and title
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Create title text
            base_name = display_names.get(plot_type, plot_type.replace('_', ' ').title())
            if offset_info and offset_info.get('offset_name'):
                slide_title = f"{base_name}: {offset_info['offset_name']}"
            else:
                slide_title = base_name
            
            # Add title text box
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_title
            title_para.font.name = "Times New Roman"
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            # Add image below title
            left = Inches(0.5)
            top = Inches(0.8)
            width = Inches(12.333)
            
            slide.shapes.add_picture(str(temp_png), left, top, width=width)
        
        prs.save(output_path)

    def _generate_single_plot(self, generator, plot_type: str, output_path: str, config: PlotConfig, 
                              offset_info: Optional[Dict] = None):
        """Generate a single plot type.
        
        Args:
            generator: PublicationFigureGenerator instance
            plot_type: Internal key of the plot type
            output_path: Path to save the figure
            config: PlotConfig instance
            offset_info: Optional dict with offset details (for individual offset plots)
        """
        if plot_type == 'aggregated':
            generator.generate_aggregated_plot(output_path=output_path, config=config)
        elif plot_type == 'per_offset':
            max_offsets = self.max_offsets_spin.value()
            generator.generate_per_offset_plot(
                output_path=output_path,
                config=config,
                max_offsets=max_offsets
            )
        elif plot_type == 'uncertainty':
            generator.generate_uncertainty_plot(output_path=output_path, config=config)
        elif plot_type == 'aggregated_wavelength':
            generator.generate_aggregated_wavelength_plot(output_path=output_path, config=config)
        elif plot_type == 'per_offset_wavelength':
            max_offsets = self.max_offsets_spin.value()
            generator.generate_per_offset_wavelength_plot(
                output_path=output_path,
                config=config,
                max_offsets=max_offsets
            )
        elif plot_type == 'dual_domain':
            generator.generate_dual_domain_plot(output_path=output_path, config=config)
        # Canvas Export types
        elif plot_type == 'canvas_frequency':
            generator.generate_canvas_frequency(output_path=output_path, config=config)
        elif plot_type == 'canvas_wavelength':
            generator.generate_canvas_wavelength(output_path=output_path, config=config)
        elif plot_type == 'canvas_dual':
            generator.generate_canvas_dual(output_path=output_path, config=config)
        # Source Offset Analysis types (individual offsets)
        elif plot_type == 'offset_curve_only':
            if offset_info and offset_info.get('offset_index') is not None:
                generator.generate_offset_curve_only(
                    output_path=output_path, 
                    config=config,
                    offset_index=offset_info['offset_index']
                )
            else:
                raise ValueError("offset_curve_only requires offset_info with offset_index")
        elif plot_type == 'offset_with_spectrum':
            if offset_info and offset_info.get('offset_index') is not None:
                generator.generate_offset_with_spectrum(
                    output_path=output_path, 
                    config=config,
                    offset_index=offset_info['offset_index']
                )
            else:
                raise ValueError("offset_with_spectrum requires offset_info with offset_index")
        elif plot_type == 'offset_spectrum_only':
            if offset_info and offset_info.get('offset_index') is not None:
                generator.generate_offset_spectrum_only(
                    output_path=output_path, 
                    config=config,
                    offset_index=offset_info['offset_index']
                )
            else:
                raise ValueError("offset_spectrum_only requires offset_info with offset_index")
        elif plot_type == 'offset_grid':
            # Grid uses grid layout and display mode from offset options
            rows_widget = getattr(self, 'grid_rows_spin', None)
            cols_widget = getattr(self, 'grid_cols_spin', None)
            rows = rows_widget.value() if rows_widget else None
            cols = cols_widget.value() if cols_widget else None
            
            # Determine display mode from radio buttons
            include_spectrum = False
            include_curves = True
            if hasattr(self, 'grid_mode_spectrum') and self.grid_mode_spectrum.isChecked():
                include_spectrum = True
                include_curves = False
            elif hasattr(self, 'grid_mode_both') and self.grid_mode_both.isChecked():
                include_spectrum = True
                include_curves = True
            
            # Pass spectrum data list for spectrum modes
            spectrum_data_list = generator.spectrum_data_list if include_spectrum else None
            
            generator.generate_offset_grid(
                output_path=output_path, 
                config=config, 
                rows=rows, 
                cols=cols,
                include_spectrum=include_spectrum,
                spectrum_data_list=spectrum_data_list,
                include_curves=include_curves
            )
        else:
            raise NotImplementedError(f"Plot type '{plot_type}' is not yet implemented.")
